"""
FileParser —— 多格式文件解析组件。

支持 PDF / PPTX / CSV / Excel / 图片 / 音视频。
限制：单文件 50MB，注入文本上限 4000 字符。

满足：R7, R14
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional


@dataclass
class ParsedFileRecord:
    file_id: str
    mime_type: str
    summary: str
    extracted_text: str
    preview_payload: Any = None


class FileParser:
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
    MAX_INJECT_CHARS = 4000

    # MIME 类型到 handler 的映射
    HANDLERS = {
        "application/pdf": "_parse_pdf",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "_parse_pptx",
        "text/csv": "_parse_csv",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "_parse_excel",
        "application/vnd.ms-excel": "_parse_excel",
        "image/png": "_parse_image",
        "image/jpeg": "_parse_image",
        "image/jpg": "_parse_image",
        "audio/mpeg": "_parse_audio",
        "audio/wav": "_parse_audio",
        "video/mp4": "_parse_video",
        "video/quicktime": "_parse_video",
    }

    async def parse(self, file_path: str, mime_type: str) -> ParsedFileRecord:
        """解析文件并返回 ParsedFileRecord。"""
        import uuid

        file_id = uuid.uuid4().hex[:8]

        # 检查文件大小
        try:
            file_size = os.path.getsize(file_path)
        except OSError:
            return ParsedFileRecord(
                file_id=file_id,
                mime_type=mime_type,
                summary="文件不存在或无法访问",
                extracted_text="",
            )

        if file_size > self.MAX_FILE_SIZE:
            return ParsedFileRecord(
                file_id=file_id,
                mime_type=mime_type,
                summary="file_too_large",
                extracted_text="",
            )

        handler_name = self.HANDLERS.get(mime_type)
        if handler_name is None:
            return ParsedFileRecord(
                file_id=file_id,
                mime_type=mime_type,
                summary=f"不支持的文件类型: {mime_type}",
                extracted_text="",
            )

        handler = getattr(self, handler_name, None)
        if handler is None:
            return ParsedFileRecord(
                file_id=file_id,
                mime_type=mime_type,
                summary=f"Handler 未实现: {handler_name}",
                extracted_text="",
            )

        try:
            record = await handler(file_path)
            record.file_id = file_id
            record.mime_type = mime_type

            # 截断注入文本
            if len(record.extracted_text) > self.MAX_INJECT_CHARS:
                record.extracted_text = record.extracted_text[: self.MAX_INJECT_CHARS] + "…(truncated)"

            return record
        except Exception as e:
            return ParsedFileRecord(
                file_id=file_id,
                mime_type=mime_type,
                summary=f"解析失败: {str(e)}",
                extracted_text="",
            )

    async def _parse_pdf(self, file_path: str) -> ParsedFileRecord:
        """PDF 解析：提取所有页文本 + 首页缩略图。"""
        text_parts = []
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
        except ImportError:
            text_parts.append("[pdfplumber 未安装，无法解析 PDF]")
        except Exception as e:
            text_parts.append(f"[PDF 解析错误: {e}]")

        extracted = "\n".join(text_parts)
        # 首页缩略图（try pdf2image）
        thumbnail = None
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(file_path, first_page=1, last_page=1)
            if images:
                thumbnail = images[0]  # PIL Image
        except Exception:
            pass

        return ParsedFileRecord(
            file_id="",
            mime_type="application/pdf",
            summary=f"PDF 文档，{len(text_parts)} 页",
            extracted_text=extracted,
            preview_payload={"thumbnail": thumbnail is not None, "pages": len(text_parts)},
        )

    async def _parse_pptx(self, file_path: str) -> ParsedFileRecord:
        """PPTX 解析：按幻灯片提取标题、正文、备注。"""
        slides = []
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_data = {"index": i + 1, "title": "", "text": "", "notes": ""}
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                            slide_data["title"] = text
                        else:
                            slide_data["text"] += text + "\n"
                if slide.has_notes_slide:
                    slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
                slides.append(slide_data)
        except ImportError:
            slides.append({"index": 1, "title": "", "text": "[python-pptx 未安装]",
                           "notes": ""})
        except Exception as e:
            slides.append({"index": 1, "title": "", "text": f"[PPTX 解析错误: {e}]",
                           "notes": ""})

        extracted = "\n".join(
            f"Slide {s['index']}: {s['title']}\n{s['text']}\n备注: {s['notes']}"
            for s in slides
        )
        return ParsedFileRecord(
            file_id="",
            mime_type="",
            summary=f"PPTX 演示文稿，{len(slides)} 张幻灯片",
            extracted_text=extracted,
            preview_payload={"slides": len(slides)},
        )

    async def _parse_csv(self, file_path: str) -> ParsedFileRecord:
        """CSV 解析：编码探测 + DuckDB 入库 + 前 20 行预览。"""
        import pandas as pd

        # 编码探测
        encodings = ["utf-8", "gbk", "gb18030", "latin-1"]
        df = None
        used_encoding = "utf-8"
        for enc in encodings:
            try:
                df = pd.read_csv(file_path, encoding=enc, nrows=5000)
                used_encoding = enc
                break
            except (UnicodeDecodeError, UnicodeError):
                continue

        if df is None:
            return ParsedFileRecord(
                file_id="", mime_type="text/csv",
                summary="CSV 编码探测失败",
                extracted_text="",
            )

        preview = df.head(20).to_dict(orient="records")
        return ParsedFileRecord(
            file_id="", mime_type="text/csv",
            summary=f"CSV 文件，{len(df.columns)} 列，编码: {used_encoding}",
            extracted_text=df.head(100).to_csv(index=False),
            preview_payload={"columns": list(df.columns), "rows": preview, "encoding": used_encoding},
        )

    async def _parse_excel(self, file_path: str) -> ParsedFileRecord:
        """Excel 解析：所有 sheet 的前 20 行预览。"""
        import pandas as pd

        try:
            xls = pd.ExcelFile(file_path)
            sheets = {}
            all_text = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=5000)
                sheets[sheet_name] = df.head(20).to_dict(orient="records")
                all_text.append(f"--- Sheet: {sheet_name} ---\n{df.head(100).to_csv(index=False)}")
            return ParsedFileRecord(
                file_id="", mime_type="",
                summary=f"Excel 文件，{len(xls.sheet_names)} 个工作表",
                extracted_text="\n".join(all_text),
                preview_payload={"sheets": sheets},
            )
        except Exception as e:
            return ParsedFileRecord(
                file_id="", mime_type="",
                summary=f"Excel 解析失败: {e}",
                extracted_text="",
            )

    async def _parse_image(self, file_path: str) -> ParsedFileRecord:
        """图片解析：OCR（paddleocr 优先，降级为仅返回文件名）。"""
        text = ""
        try:
            from paddleocr import PaddleOCR
            ocr = PaddleOCR(lang="ch")
            result = ocr.ocr(file_path)
            if result and result[0]:
                text = "\n".join(line[1][0] for line in result[0])
        except ImportError:
            pass
        except Exception:
            pass

        if not text:
            text = f"[图片文件: {os.path.basename(file_path)}，OCR 不可用]"

        return ParsedFileRecord(
            file_id="", mime_type="",
            summary=f"图片文件: {os.path.basename(file_path)}",
            extracted_text=text,
            preview_payload={"ocr_available": bool(text and not text.startswith("[图片文件"))},
        )

    async def _parse_audio(self, file_path: str) -> ParsedFileRecord:
        """音频解析：Whisper 转写，降级为仅返回文件名与时长。"""
        text = ""
        duration = ""
        try:
            import whisper
            model = whisper.load_model("base")
            result = model.transcribe(file_path)
            text = result.get("text", "")
        except ImportError:
            pass
        except Exception:
            pass

        if not text:
            text = f"[音频文件: {os.path.basename(file_path)}，转写不可用]"

        return ParsedFileRecord(
            file_id="", mime_type="",
            summary=f"音频文件: {os.path.basename(file_path)}",
            extracted_text=text,
            preview_payload={"duration": duration},
        )

    async def _parse_video(self, file_path: str) -> ParsedFileRecord:
        """视频解析：Whisper 转写音频轨，降级为仅返回文件名。"""
        # 视频转写：优先尝试提取音频再 transcribe
        text = ""
        try:
            import whisper
            model = whisper.load_model("base")
            result = model.transcribe(file_path)
            text = result.get("text", "")
        except ImportError:
            pass
        except Exception:
            pass

        if not text:
            text = f"[视频文件: {os.path.basename(file_path)}，转写不可用]"

        return ParsedFileRecord(
            file_id="", mime_type="",
            summary=f"视频文件: {os.path.basename(file_path)}",
            extracted_text=text,
        )


# 单例
_file_parser: Optional[FileParser] = None


def get_file_parser() -> FileParser:
    global _file_parser
    if _file_parser is None:
        _file_parser = FileParser()
    return _file_parser
